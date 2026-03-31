from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Color Palette ─────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x0D, 0x1B, 0x2A)   # near-black navy
ACCENT      = RGBColor(0x4A, 0xC9, 0xFF)   # cyan-blue
ACCENT2     = RGBColor(0x7C, 0x3A, 0xFF)   # purple
GREEN       = RGBColor(0x2E, 0xCC, 0x71)   # success green
YELLOW      = RGBColor(0xFF, 0xD7, 0x00)   # gold/yellow
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xB0, 0xBE, 0xC5)
CARD_BG     = RGBColor(0x16, 0x2A, 0x3E)   # slightly lighter than dark bg
RED         = RGBColor(0xFF, 0x5C, 0x5C)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, alpha=None, line_color=None, line_width=None, radius=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, radius_pt=6, line_color=None, line_width=1):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    # Set corner radius
    adj = shape.adjustments
    if len(adj) > 0:
        adj[0] = 0.05  # small rounding
    return shape

def txb(slide, text, x, y, w, h,
        font_size=18, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT,
        font_name="Segoe UI", word_wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = word_wrap
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

def txb_multi(slide, lines, x, y, w, h,
              font_size=14, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, spacing_after=0, font_name="Segoe UI"):
    """lines = list of (text, bold, color, size) or just strings"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, c, s = item, bold, color, font_size
        else:
            txt = item.get("text","")
            b   = item.get("bold", bold)
            c   = item.get("color", color)
            s   = item.get("size", font_size)
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if spacing_after:
            p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(s)
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = font_name
    return tb

def bg(slide, color=DARK_BG):
    r = add_rect(slide, 0, 0, 13.33, 7.5, fill_color=color)
    # send to back
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)

def accent_bar(slide, x=0, y=7.1, w=13.33, h=0.08, color=ACCENT):
    add_rect(slide, x, y, w, h, fill_color=color)

def slide_number(slide, n, total=12):
    txb(slide, f"{n} / {total}", 12.5, 7.15, 0.7, 0.25,
        font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def header_tag(slide, tag_text, color=ACCENT):
    """Small colored pill label top-left"""
    add_rounded_rect(slide, 0.4, 0.25, 2.2, 0.38, fill_color=color)
    txb(slide, tag_text, 0.44, 0.25, 2.1, 0.38,
        font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

def divider(slide, y, color=ACCENT, x=0.4, w=12.53):
    add_rect(slide, x, y, w, 0.03, fill_color=color)

def icon_bullet(slide, items, x, y, w, line_h=0.38, dot_color=ACCENT,
                font_size=13, text_color=WHITE, bold=False):
    """items = list of strings"""
    cy = y
    for item in items:
        add_rect(slide, x, cy+0.12, 0.1, 0.1, fill_color=dot_color)
        txb(slide, item, x+0.18, cy, w-0.18, line_h,
            font_size=font_size, color=text_color, bold=bold)
        cy += line_h

def stat_card(slide, label, value, unit, x, y, w=2.8, h=1.5,
              accent=ACCENT, bg_color=CARD_BG):
    add_rounded_rect(slide, x, y, w, h, fill_color=bg_color)
    add_rect(slide, x, y, 0.05, h, fill_color=accent)
    txb(slide, value, x+0.15, y+0.12, w-0.2, 0.75,
        font_size=36, bold=True, color=accent, align=PP_ALIGN.LEFT)
    txb(slide, unit, x+0.15, y+0.72, w-0.2, 0.4,
        font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
    txb(slide, label, x+0.15, y+1.1, w-0.2, 0.35,
        font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# =============================================================================
# SLIDE 1 - TITLE
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)

# gradient overlay bands
add_rect(s, 0, 0, 13.33, 3.8, fill_color=RGBColor(0x0A,0x14,0x20))
add_rect(s, 0, 0, 0.5, 7.5, fill_color=ACCENT2)
add_rect(s, 0, 0, 0.08, 7.5, fill_color=ACCENT)

# big title
txb(s, "StudiBudd", 1.0, 1.4, 9, 1.6,
    font_size=80, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# cyan underline
add_rect(s, 1.0, 2.9, 5.5, 0.07, fill_color=ACCENT)

# tagline
txb(s, "Gamified Academic Productivity for the Modern College Student",
    1.0, 3.1, 10, 0.6,
    font_size=22, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

# decorative circles (top right)
for r_size, alpha_col in [(2.5, RGBColor(0x1A,0x2E,0x42)),
                           (1.7, RGBColor(0x16,0x26,0x3A)),
                           (1.0, RGBColor(0x4A,0xC9,0xFF))]:
    circ = s.shapes.add_shape(9, Inches(11.5), Inches(0.2),
                               Inches(r_size), Inches(r_size))
    circ.fill.solid()
    circ.fill.fore_color.rgb = alpha_col
    circ.line.fill.background()

# bottom info
add_rect(s, 0, 6.6, 13.33, 0.9, fill_color=RGBColor(0x08,0x10,0x18))
txb(s, "Entrepreneurial Management  |  University of New Haven  |  2026",
    1.0, 6.68, 11, 0.4, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

accent_bar(s, color=ACCENT2)

# =============================================================================
# SLIDE 2 - THE PROBLEM
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 2)
header_tag(s, "CURIOSITY  |  01")

txb(s, "The Problem", 0.4, 0.7, 8, 0.75,
    font_size=42, bold=True, color=WHITE)
txb(s, "Students are overwhelmed — and their tools aren't built for them.",
    0.4, 1.38, 10, 0.45, font_size=16, italic=True, color=LIGHT_GRAY)

divider(s, 1.88)

# left column - stat cards
stat_card(s, "Report constant academic stress",      "40%", "of college students",
          0.4,  2.1, 3.1, 1.4, accent=ACCENT)
stat_card(s, "Apps used to manage coursework",       "5+",  "separate apps avg.",
          0.4,  3.65, 3.1, 1.4, accent=ACCENT2)
stat_card(s, "Feel so depressed it's hard to focus", "44%", "of students (ACHA 2023)",
          0.4,  5.2,  3.1, 1.4, accent=RED)

# right column - the gap
add_rounded_rect(s, 3.9, 2.1, 9.0, 4.7, fill_color=CARD_BG)
add_rect(s, 3.9, 2.1, 9.0, 0.05, fill_color=ACCENT)

txb(s, "The Gap in the Market", 4.1, 2.2, 8.5, 0.5,
    font_size=18, bold=True, color=ACCENT)

gaps = [
    "Notion / Google Calendar  -  Generic; built for professionals, not students",
    "Canvas LMS  -  Assignment tracking only; no motivation layer",
    "Forest / Focus apps  -  Motivation only; no academic organization",
    "MyStudyLife  -  Manual entry only; no LMS integration",
]
icon_bullet(s, gaps, 4.1, 2.85, 8.5, line_h=0.5, dot_color=RED, font_size=13)

divider(s, 5.1, color=ACCENT2, x=4.1, w=8.5)

txb(s, "The Insight:", 4.1, 5.22, 2.0, 0.4, font_size=14, bold=True, color=YELLOW)
txb(s, "Students don't just need a task list — they need a reason to open the app every single day. Habit formation requires emotional engagement, not just functionality.",
    4.1, 5.22, 8.6, 0.9, font_size=13, color=WHITE)

txb(s, u"Sources: \u00b940% & 44% \u2014 ACHA-NCHA Spring 2023 Reference Group Executive Summary (N=54,497)  |  \u00b95+ apps \u2014 internal survey of 42 UNH students, Spring 2026",
    0.4, 6.78, 12.5, 0.28, font_size=8, italic=True, color=LIGHT_GRAY)

# =============================================================================
# SLIDE 3 - MARKET RESEARCH
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 3)
header_tag(s, "CURIOSITY  |  02")

txb(s, "Market Research & Sizing", 0.4, 0.7, 9, 0.7, font_size=42, bold=True, color=WHITE)
txb(s, "Extensive data-driven research with validated assumptions.",
    0.4, 1.35, 9, 0.4, font_size=15, italic=True, color=LIGHT_GRAY)
divider(s, 1.82)

# TAM / SAM / SOM - three cards
for i, (label, val, sub, col) in enumerate([
    ("TAM", "$5.2B", "Total Addressable Market\n20M US college students\nEdTech productivity (12% CAGR)", ACCENT),
    ("SAM", "5M",    "Serviceable Addressable Market\nStudents using digital study tools\nAt 4-year universities", ACCENT2),
    ("SOM", "75K",   "Serviceable Obtainable Market\nTarget users by Year 3\nEntry via UNH + CT/NY schools", GREEN),
]):
    cx = 0.4 + i * 4.3
    add_rounded_rect(s, cx, 2.1, 4.0, 2.8, fill_color=CARD_BG)
    add_rect(s, cx, 2.1, 4.0, 0.07, fill_color=col)
    txb(s, label, cx+0.2, 2.2, 3.5, 0.55, font_size=28, bold=True, color=col)
    txb(s, val,   cx+0.2, 2.72, 3.5, 0.75, font_size=38, bold=True, color=WHITE)
    txb(s, sub,   cx+0.2, 3.45, 3.5, 1.35, font_size=11, color=LIGHT_GRAY)

# validated assumptions table
add_rounded_rect(s, 0.4, 5.1, 12.5, 2.1, fill_color=CARD_BG)
txb(s, "Validated Assumptions", 0.6, 5.18, 5, 0.4, font_size=14, bold=True, color=ACCENT)

headers = ["Assumption", "Validation", "Confidence"]
col_x   = [0.6, 5.5, 10.5]
col_w   = [4.7, 4.8, 2.2]
for hx, hw, ht in zip(col_x, col_w, headers):
    txb(s, ht, hx, 5.55, hw, 0.3, font_size=11, bold=True, color=ACCENT2)

rows = [
    ("$4.99/mo price point acceptable",        "Spotify Student: $5.99/mo (spotify.com/us/student)", "HIGH"),
    ("Students check planner apps daily",      "Forest: 10M+ DAU (App Store, 2024); habit loop lit.","HIGH"),
    ("Canvas at 30%+ of US universities",      "Instructure 2023 Annual Report, p.4",                "HIGH"),
    ("Universities spend $1.2B on EdTech",     "HolonIQ Global EdTech Intelligence Report 2024",     "HIGH"),
]
for ri, row in enumerate(rows):
    ry = 5.88 + ri * 0.3
    for ci, (cx, cw, cell) in enumerate(zip(col_x, col_w, row)):
        col_c = GREEN if cell == "HIGH" else WHITE
        txb(s, cell, cx, ry, cw, 0.28, font_size=10,
            color=col_c, bold=(ci==2))

txb(s, u"Sources: \u00b9NCES Digest of Education Statistics 2023 (20M students)  |  \u00b2Grand View Research EdTech Market Report 2024 ($5.2B / 12% CAGR)  |  \u00b3Instructure 2023 Annual Report  |  \u2074HolonIQ 2024",
    0.4, 6.78, 12.5, 0.28, font_size=8, italic=True, color=LIGHT_GRAY)

# =============================================================================
# SLIDE 4 - BUSINESS MODEL CANVAS
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 4)
header_tag(s, "CURIOSITY  |  03")

txb(s, "Business Model Canvas", 0.4, 0.7, 9, 0.7, font_size=42, bold=True, color=WHITE)
divider(s, 1.45)

# BMC grid
# Row 1: Key Partners | Key Activities | Value Props | Cust Relationships | Segments
# Row 2: Key Partners | Key Resources  | Value Props | Cust Relationships | Segments
# Row 3: Cost Structure (wide)         | Revenue Streams (wide)

bmc_data = [
    # (label, content, x, y, w, h, accent_color)
    ("KEY PARTNERS",
     "University IT depts\nCanvas / Instructure\nVercel + Supabase\nApp stores",
     0.25, 1.65, 2.4, 2.4, ACCENT2),
    ("KEY ACTIVITIES",
     "App dev & maintenance\nCampus ambassador mgmt\nUniversity sales cycle\nFeature iteration",
     2.73, 1.65, 2.4, 1.15, ACCENT),
    ("VALUE PROPOSITIONS",
     "One app replaces 5\nReal Canvas sync — no manual entry\nCreature gamification loop\nCustomizable aesthetic",
     5.21, 1.65, 2.9, 2.4, YELLOW),
    ("CUSTOMER REL.",
     "Gamification & streaks\nAchievement badges\nEmail support",
     8.19, 1.65, 2.4, 1.15, ACCENT2),
    ("CUSTOMER SEGMENTS",
     "PRIMARY:\nCollege students 18-24\n\nSECONDARY:\nUniversities (B2B2C)",
     10.67, 1.65, 2.4, 2.4, ACCENT),
    ("KEY RESOURCES",
     "Next.js codebase\nCanvas API integration\nTeam expertise\nUser data",
     2.73, 2.88, 2.4, 1.17, ACCENT),
    ("CUSTOMER REL. (2)",
     "Campus ambassadors\nProfessor endorsements\nWord of mouth",
     8.19, 2.88, 2.4, 1.17, ACCENT2),
    ("COST STRUCTURE",
     "Hosting: $540/yr  |  Domain/tools: $212/yr  |  Marketing: $500/yr (Y1)  |  Dev (opp. cost): $10,000",
     0.25, 4.13, 6.3, 1.2, RED),
    ("REVENUE STREAMS",
     "Freemium: $4.99/mo premium  |  University licensing: $2,000-5,000/yr per institution",
     6.63, 4.13, 6.4, 1.2, GREEN),
]

for (lbl, content, bx, by, bw, bh, bcol) in bmc_data:
    if lbl == "VALUE PROPOSITIONS":
        bh_use = bh
    else:
        bh_use = bh
    add_rounded_rect(s, bx, by, bw, bh_use, fill_color=CARD_BG, line_color=bcol, line_width=1.2)
    add_rect(s, bx, by, bw, 0.05, fill_color=bcol)
    txb(s, lbl, bx+0.1, by+0.06, bw-0.15, 0.28, font_size=8, bold=True, color=bcol)
    txb(s, content, bx+0.1, by+0.32, bw-0.15, bh_use-0.35,
        font_size=9, color=WHITE)

txb(s, "The Canvas integration is not just a feature — it is where every business module converges.",
    0.4, 5.45, 12.5, 0.38, font_size=12, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 5 - COST & REVENUE MODEL
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 5)
header_tag(s, "CONNECTIONS  |  04")

txb(s, "Cost & Revenue Model", 0.4, 0.7, 9, 0.7, font_size=42, bold=True, color=WHITE)
txb(s, "Lean cost structure — high-margin, cloud-native business.",
    0.4, 1.35, 9, 0.4, font_size=15, italic=True, color=LIGHT_GRAY)
divider(s, 1.82)

# ── Revenue table (left) ──────────────────────────────────────────────────────
add_rounded_rect(s, 0.4, 2.05, 7.8, 4.8, fill_color=CARD_BG)
add_rect(s, 0.4, 2.05, 7.8, 0.06, fill_color=GREEN)
txb(s, "Revenue Projections", 0.6, 2.14, 5, 0.4, font_size=15, bold=True, color=GREEN)

rev_headers = ["Stream", "Year 1", "Year 2", "Year 3"]
rev_cx = [0.6, 3.5, 5.15, 6.8]
rev_cw = [2.8, 1.55, 1.55, 1.5]
for hx, hw, ht in zip(rev_cx, rev_cw, rev_headers):
    txb(s, ht, hx, 2.58, hw, 0.32, font_size=11, bold=True, color=ACCENT)

rev_rows = [
    ("Premium Users",          "200",      "500",      "1,200"),
    ("Premium Revenue",        "$11,976",  "$29,940",  "$71,856"),
    ("University Partners",    "1",        "3",        "8"),
    ("University Revenue",     "$2,000",   "$6,000",   "$16,000"),
    ("", "", "", ""),
    ("TOTAL REVENUE",          "$13,976",  "$35,940",  "$87,856"),
    ("Operating Costs",        "($1,252)", "($2,452)", "($4,492)"),
    ("NET CASH FLOW",          "$12,724",  "$33,488",  "$83,364"),
]

for ri, row in enumerate(rev_rows):
    ry = 2.95 + ri * 0.47
    is_total = row[0] in ("TOTAL REVENUE", "NET CASH FLOW")
    if is_total:
        add_rect(s, 0.45, ry-0.05, 7.7, 0.42, fill_color=RGBColor(0x1E,0x38,0x50))
    for ci, (cx, cw, cell) in enumerate(zip(rev_cx, rev_cw, row)):
        c = GREEN if (is_total and ci > 0) else (WHITE if ci == 0 else LIGHT_GRAY)
        b = is_total
        txb(s, cell, cx, ry, cw, 0.38, font_size=11, bold=b, color=c)

# ── Cost breakdown (right) ────────────────────────────────────────────────────
add_rounded_rect(s, 8.55, 2.05, 4.4, 4.8, fill_color=CARD_BG)
add_rect(s, 8.55, 2.05, 4.4, 0.06, fill_color=RED)
txb(s, "Cost Structure (Year 1)", 8.72, 2.14, 4.0, 0.4, font_size=14, bold=True, color=RED)

cost_items = [
    ("Hosting (Vercel + Supabase)", "$540"),
    ("Domain & misc tools",         "$212"),
    ("Marketing & campus outreach", "$500"),
    ("Customer support tools",      "$0"),
    ("TOTAL OPERATING",             "$1,252"),
    ("", ""),
    ("Initial Investment",          "$10,000"),
    ("(dev time + setup)",          ""),
]
for ri, (label, val) in enumerate(cost_items):
    ry = 2.65 + ri * 0.52
    is_total = label.startswith("TOTAL") or label.startswith("Initial")
    if is_total:
        add_rect(s, 8.58, ry-0.04, 4.35, 0.44, fill_color=RGBColor(0x1E,0x38,0x50))
    col_l = WHITE if not is_total else YELLOW
    col_v = RED if is_total else ACCENT
    txb(s, label, 8.72, ry, 2.6, 0.38, font_size=10, bold=is_total, color=col_l)
    txb(s, val,   11.1, ry, 1.7, 0.38, font_size=11, bold=is_total, color=col_v,
        align=PP_ALIGN.RIGHT)

txb(s, u"Note: All figures are internal projections. Hosting costs from Vercel/Supabase pricing pages (2026). User growth assumptions based on comparable EdTech freemium benchmarks (Duolingo S-1: 4-8% free-to-paid conversion). University pricing benchmarked against Canvas LTI app market ($1K-$10K/yr). Initial investment = ~400 dev hrs \u00d7 $25/hr opportunity cost.",
    0.4, 6.78, 12.5, 0.28, font_size=8, italic=True, color=LIGHT_GRAY)

# =============================================================================
# SLIDE 6 - FINANCIAL ANALYSIS
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 6)
header_tag(s, "CONNECTIONS  |  05")

txb(s, "Financial Analysis", 0.4, 0.7, 9, 0.7, font_size=42, bold=True, color=WHITE)
txb(s, "NPV  |  IRR  |  Payback Period — all metrics confirm viability.",
    0.4, 1.35, 9, 0.4, font_size=15, italic=True, color=LIGHT_GRAY)
divider(s, 1.82)

# Three big KPI cards
stat_card(s, "Net Present Value (10% discount)",  "$91,845", "NPV — well above $0",
          0.4,  2.1, 3.85, 1.65, accent=GREEN)
stat_card(s, "Internal Rate of Return",           "254%",    "IRR — vs 10% hurdle rate",
          4.4,  2.1, 3.85, 1.65, accent=ACCENT)
stat_card(s, "Investment Payback Period",         "~10 mo",  "Payback within Year 1",
          8.4,  2.1, 4.5,  1.65, accent=YELLOW)

# Cash flow table
add_rounded_rect(s, 0.4, 3.95, 12.5, 2.9, fill_color=CARD_BG)
add_rect(s, 0.4, 3.95, 12.5, 0.05, fill_color=ACCENT2)
txb(s, "Discounted Cash Flow Analysis  (Discount Rate: 10%)", 0.6, 4.03, 9, 0.4,
    font_size=13, bold=True, color=ACCENT2)

cf_headers = ["Year", "Net Cash Flow", "Discount Factor", "Present Value", "Cumulative PV"]
cf_cx = [0.6, 2.5, 5.0, 7.5, 10.2]
cf_cw = [1.8, 2.4, 2.4, 2.6, 2.5]

for hx, hw, ht in zip(cf_cx, cf_cw, cf_headers):
    txb(s, ht, hx, 4.48, hw, 0.32, font_size=11, bold=True, color=ACCENT)

cf_rows = [
    ("0 (Investment)", "-$10,000",  "1.000", "-$10,000",  "-$10,000"),
    ("1",              "$12,724",   "0.909",  "$11,567",   "$1,567"),
    ("2",              "$33,488",   "0.826",  "$27,676",  "$29,243"),
    ("3",              "$83,364",   "0.751",  "$62,606",  "$91,849"),
]
for ri, row in enumerate(cf_rows):
    ry = 4.86 + ri * 0.48
    is_last = ri == 3
    if is_last:
        add_rect(s, 0.45, ry-0.04, 12.4, 0.44, fill_color=RGBColor(0x1E,0x38,0x50))
    for ci, (cx, cw, cell) in enumerate(zip(cf_cx, cf_cw, row)):
        color = GREEN if (is_last and ci >= 3) else (WHITE if ci == 0 else LIGHT_GRAY)
        txb(s, cell, cx, ry, cw, 0.38, font_size=11, bold=is_last, color=color)

# interpretation note
txb(s, "Decision: NPV > 0, IRR >> hurdle rate, Payback < 1 year.  All three criteria independently justify investment.",
    0.4, 6.72, 12.5, 0.3, font_size=11, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

txb(s, u"Methodology: NPV = \u03a3[CFt / (1+r)^t] \u2212 I\u2080  |  IRR solved via =IRR() formula  |  Discount rate 10% = standard VC hurdle (Damodaran, NYU 2024)  |  Cash flows derived from internal revenue model (Slide 5)",
    0.4, 6.98, 12.5, 0.28, font_size=8, italic=True, color=LIGHT_GRAY)

# =============================================================================
# SLIDE 7 - EM MODULE INTEGRATION
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 7)
header_tag(s, "CONNECTIONS  |  06")

txb(s, "Integration of EM Modules", 0.4, 0.7, 9, 0.7, font_size=42, bold=True, color=WHITE)
txb(s, "Systems thinking — every decision reinforces the others.",
    0.4, 1.35, 10, 0.4, font_size=15, italic=True, color=LIGHT_GRAY)
divider(s, 1.82)

# ── helper: draw a straight connector with an arrowhead at the END ──────────
from pptx.oxml.ns import qn as _qn

def add_arrow(slide, x1, y1, x2, y2, color, width_pt=2.0, dashed=False):
    """Draws a straight line from (x1,y1) to (x2,y2) with an arrowhead at the end."""
    conn = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    if dashed:
        conn.line.dash_style = 4   # DASH
    # Inject arrowhead via XML on the a:ln element
    sp_pr = conn._element.spPr
    ln = sp_pr.find(_qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(sp_pr, _qn('a:ln'))
    tail = etree.SubElement(ln, _qn('a:tailEnd'))
    tail.set('type', 'arrow')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return conn

def add_connector_label(slide, label, cx, cy, color, font_size=9):
    """Small pill label placed along a connector."""
    w, h = 1.9, 0.28
    add_rounded_rect(slide, cx - w/2, cy - h/2, w, h,
                     fill_color=DARK_BG, line_color=color, line_width=0.8)
    txb(slide, label, cx - w/2 + 0.05, cy - h/2, w - 0.05, h,
        font_size=font_size, color=color, align=PP_ALIGN.CENTER, bold=True)

# ── Layout ───────────────────────────────────────────────────────────────────
# Hub (ellipse) — center at (6.65, 3.80)
HUB_X, HUB_Y, HUB_W, HUB_H = 5.55, 2.95, 2.2, 1.7
HUB_CX = HUB_X + HUB_W / 2   # 6.65
HUB_CY = HUB_Y + HUB_H / 2   # 3.80

hub = s.shapes.add_shape(9,   # ROUNDED_RECTANGLE
    Inches(HUB_X), Inches(HUB_Y), Inches(HUB_W), Inches(HUB_H))
hub.fill.solid()
hub.fill.fore_color.rgb = ACCENT
hub.line.color.rgb = WHITE
hub.line.width = Pt(1.5)
txb(s, "CANVAS\nINTEGRATION\nHUB", HUB_X, HUB_Y, HUB_W, HUB_H,
    font_size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# ── 5 module boxes  (x, y, w=3.0, h=1.55) ───────────────────────────────────
MOD_W, MOD_H = 3.0, 1.55
modules = [
    # label lines, x,   y,     color,   connector: from_xy -> to_xy, label_xy, label_text
    ("MARKETING\nCampus ambassadors\nFreemium funnel\nLow-cost acquisition",
     0.25, 1.85, ACCENT2,
     (3.25, 2.63), (HUB_X, 3.40), (4.37, 3.02), "drives users"),

    ("OPERATIONS\nCloud-native deploy\nAPI abstraction layer\nCanvas sync engine",
     10.05, 1.85, GREEN,
     (10.05, 2.63), (HUB_X + HUB_W, 3.40), (8.93, 3.02), "enables sync"),

    ("STRATEGY\nFirst-mover EdTech\nUniversity B2B2C\nCanvas as moat",
     0.25, 4.85, YELLOW,
     (3.25, 5.63), (HUB_X, 4.25), (4.37, 4.94), "shapes pricing"),

    ("FINANCE\nLean cost structure\nHigh LTV / low CAC\nNPV-positive Yr 1",
     10.05, 4.85, RED,
     (10.05, 5.63), (HUB_X + HUB_W, 4.25), (8.93, 4.94), "validates ROI"),

    ("INNOVATION\nGamification loop\nAI-ready architecture\nSemester refresh",
     5.17, 6.05, RGBColor(0xFF,0x80,0x00),
     (HUB_CX, 6.05), (HUB_CX, HUB_Y + HUB_H), (HUB_CX + 1.1, 5.55), "fuels retention"),
]

for (txt, mx, my, mc, from_pt, to_pt, lbl_pt, lbl_txt) in modules:
    # draw connector FIRST (so box renders on top)
    add_arrow(s, from_pt[0], from_pt[1], to_pt[0], to_pt[1], mc, width_pt=2.2)
    add_connector_label(s, lbl_txt, lbl_pt[0], lbl_pt[1], mc)

for (txt, mx, my, mc, from_pt, to_pt, lbl_pt, lbl_txt) in modules:
    # draw boxes on top of connectors
    add_rounded_rect(s, mx, my, MOD_W, MOD_H, fill_color=CARD_BG, line_color=mc, line_width=1.8)
    add_rect(s, mx, my, MOD_W, 0.06, fill_color=mc)
    lines = txt.split("\n")
    txb(s, lines[0], mx+0.12, my+0.08, MOD_W-0.18, 0.32,
        font_size=12, bold=True, color=mc)
    txb(s, "\n".join(lines[1:]), mx+0.12, my+0.42, MOD_W-0.18, MOD_H-0.48,
        font_size=10, color=LIGHT_GRAY)

# redraw hub on top of everything
hub2 = s.shapes.add_shape(9,
    Inches(HUB_X), Inches(HUB_Y), Inches(HUB_W), Inches(HUB_H))
hub2.fill.solid()
hub2.fill.fore_color.rgb = ACCENT
hub2.line.color.rgb = WHITE
hub2.line.width = Pt(1.5)
txb(s, "CANVAS\nINTEGRATION\nHUB", HUB_X, HUB_Y + 0.2, HUB_W, HUB_H - 0.2,
    font_size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

txb(s, "Every module connects to the Canvas hub — acquisition, operations, strategy, finance, and innovation all reinforce each other.",
    0.4, 7.05, 12.5, 0.38, font_size=11, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 8 - STRATEGIC RECOMMENDATION
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 8)
header_tag(s, "CREATING VALUE  |  07", color=ACCENT2)

txb(s, "Strategic Recommendation", 0.4, 0.7, 10, 0.7, font_size=40, bold=True, color=WHITE)
txb(s, "University-partnership-led growth — starting at UNH, scaling nationally.",
    0.4, 1.35, 10, 0.4, font_size=15, italic=True, color=LIGHT_GRAY)
divider(s, 1.82)

# three phases
phases = [
    ("Phase 1", "Now - Month 6", "Prove It at UNH",
     "Launch to UNH student body\nTarget: 200 premium users\n1 university contract\nCollect NPS + testimonials",
     ACCENT, 0.4),
    ("Phase 2", "Month 6 - 18", "Regional Expansion",
     "Pitch 10 CT/NY universities\nCampus ambassador program\nTarget: 500 users\n3 university contracts",
     ACCENT2, 4.5),
    ("Phase 3", "Month 18 - Year 3", "National Scale",
     "App store listing\n'Built by students' PR push\nTarget: 1,200 users\n8+ university contracts",
     GREEN, 8.6),
]

for (ph, timeline, title, content, col, px) in phases:
    add_rounded_rect(s, px, 2.1, 3.85, 3.6, fill_color=CARD_BG)
    add_rect(s, px, 2.1, 3.85, 0.06, fill_color=col)
    txb(s, ph, px+0.15, 2.18, 2, 0.35, font_size=11, bold=True, color=col)
    txb(s, timeline, px+0.15, 2.52, 3.5, 0.3, font_size=10, color=LIGHT_GRAY)
    txb(s, title, px+0.15, 2.85, 3.6, 0.45, font_size=16, bold=True, color=WHITE)
    txb(s, content, px+0.15, 3.38, 3.6, 2.2, font_size=12, color=LIGHT_GRAY)

# pricing rationale
add_rounded_rect(s, 0.4, 5.9, 12.5, 1.0, fill_color=CARD_BG)
add_rect(s, 0.4, 5.9, 0.06, 1.0, fill_color=YELLOW)
txb(s, "Pricing Rationale:", 0.6, 5.98, 2.5, 0.35, font_size=13, bold=True, color=YELLOW)
txb(s, "$4.99/mo premium — anchored below Spotify Student ($5.99), above 'too cheap to trust'.  "
       "University licenses at $2,000-$5,000/yr priced at <10% of typical EdTech contract value.",
    3.0, 5.98, 9.7, 0.75, font_size=12, color=WHITE)

# =============================================================================
# SLIDE 9 - VALUE PROPOSITION & STAKEHOLDER IMPACT
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 9)
header_tag(s, "CREATING VALUE  |  08", color=ACCENT2)

txb(s, "Value Proposition & Stakeholder Impact", 0.4, 0.7, 12, 0.7, font_size=36, bold=True, color=WHITE)
divider(s, 1.45)

# Three stakeholder columns
stakeholders = [
    ("STUDENTS", ACCENT,
     ["One app replaces 5",
      "Zero manual data entry via Canvas sync",
      "Creature system creates emotional engagement",
      "Course-specific deadlines with urgency flags",
      "Fully customizable aesthetic — it feels personal",
      "Free tier eliminates adoption barrier"]),
    ("UNIVERSITIES", ACCENT2,
     ["Student retention tool — reduce dropout risk",
      "No IT overhead — web app, zero install",
      "Supports Title IV reporting & success metrics",
      "Differentiator in student recruitment",
      "Data on engagement patterns (aggregate/anon)",
      "Partners in a growing EdTech ecosystem"]),
    ("SOCIETY", GREEN,
     ["Addresses 44% depression rate in students",
      "Democratizes academic org. tools",
      "Supports first-gen students w/o coaches",
      "Built by students — authentic EdTech",
      "Reduces academic burnout system-wide",
      "Scales social impact with zero marginal cost"]),
]

for i, (name, col, items) in enumerate(stakeholders):
    cx = 0.4 + i * 4.28
    add_rounded_rect(s, cx, 1.65, 4.0, 4.7, fill_color=CARD_BG)
    add_rect(s, cx, 1.65, 4.0, 0.07, fill_color=col)

    # circle icon
    circ = s.shapes.add_shape(9, Inches(cx+1.6), Inches(1.75), Inches(0.75), Inches(0.75))
    circ.fill.solid()
    circ.fill.fore_color.rgb = col
    circ.line.fill.background()

    txb(s, name, cx+0.15, 2.55, 3.7, 0.42, font_size=16, bold=True, color=col, align=PP_ALIGN.CENTER)

    for ji, item in enumerate(items):
        iy = 3.08 + ji * 0.52
        add_rect(s, cx+0.22, iy+0.14, 0.08, 0.08, fill_color=col)
        txb(s, item, cx+0.4, iy, 3.5, 0.42, font_size=11, color=WHITE)

# value statement
add_rounded_rect(s, 0.4, 6.5, 12.5, 0.75, fill_color=RGBColor(0x08, 0x1C, 0x2E))
add_rect(s, 0.4, 6.5, 0.06, 0.75, fill_color=ACCENT2)
txb(s,
    '"StudiBudd creates academic momentum for students, measurable retention metrics for universities, '
    'and a less anxious generation of graduates for society."',
    0.6, 6.55, 12.2, 0.45,
    font_size=12, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(s, u"Sources: 44% depression stat \u2014 ACHA-NCHA Spring 2023 (N=54,497)  |  Title IV retention metrics \u2014 IPEDS Graduation Rate Survey 2023  |  First-gen students \u2014 NCES 2023 First-Generation Student Report",
    0.4, 7.0, 12.5, 0.25, font_size=8, italic=True, color=LIGHT_GRAY)

# =============================================================================
# SLIDE 10 - RISK & SUSTAINABILITY
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 10)
header_tag(s, "CREATING VALUE  |  09", color=ACCENT2)

txb(s, "Risk, Sustainability & External Factors", 0.4, 0.7, 12, 0.7,
    font_size=36, bold=True, color=WHITE)
divider(s, 1.45)

# Risk table
add_rounded_rect(s, 0.4, 1.65, 8.5, 3.8, fill_color=CARD_BG)
add_rect(s, 0.4, 1.65, 8.5, 0.05, fill_color=RED)
txb(s, "Risk Register", 0.6, 1.72, 5, 0.38, font_size=14, bold=True, color=RED)

rh = ["Risk", "Prob", "Impact", "Mitigation"]
rhx = [0.6, 3.55, 4.55, 5.55]
rhw = [2.85, 0.9, 0.9, 3.2]
for hx, hw, ht in zip(rhx, rhw, rh):
    txb(s, ht, hx, 2.12, hw, 0.3, font_size=10, bold=True, color=ACCENT)

risk_rows = [
    ("Low user adoption",          "MED",  "HIGH", "Campus ambassadors + professor endorsements"),
    ("Canvas API policy change",   "LOW",  "HIGH", "Abstraction layer; monitor Instructure updates"),
    ("Competitor copies features", "HIGH", "MED",  "First-mover brand + creature emotional attachment"),
    ("FERPA / data privacy",       "LOW",  "HIGH", "No PII stored; Canvas OAuth only; privacy-by-design"),
    ("Semester churn",             "MED",  "MED",  "Creature evolution tied to new courses each term"),
]
prob_colors = {"HIGH": RED, "MED": YELLOW, "LOW": GREEN}
for ri, (r, prob, imp, mit) in enumerate(risk_rows):
    ry = 2.5 + ri * 0.53
    txb(s, r,    0.6,  ry, 2.85, 0.42, font_size=10, color=WHITE)
    txb(s, prob, 3.55, ry, 0.85, 0.42, font_size=10, bold=True, color=prob_colors[prob])
    txb(s, imp,  4.55, ry, 0.85, 0.42, font_size=10, bold=True, color=prob_colors[imp])
    txb(s, mit,  5.55, ry, 3.1,  0.42, font_size=9,  color=LIGHT_GRAY)

# Tailwinds / Headwinds
add_rounded_rect(s, 9.1, 1.65, 3.85, 1.8, fill_color=CARD_BG)
add_rect(s, 9.1, 1.65, 3.85, 0.05, fill_color=GREEN)
txb(s, "TAILWINDS", 9.25, 1.72, 3.5, 0.3, font_size=11, bold=True, color=GREEN)
tw = ["EdTech market +12% CAGR", "Canvas at 40M+ users", "Post-COVID digital habits",
      "University retention pressure"]
for i, t in enumerate(tw):
    txb(s, f"+ {t}", 9.25, 2.1+i*0.33, 3.5, 0.3, font_size=10, color=WHITE)

add_rounded_rect(s, 9.1, 3.6, 3.85, 1.85, fill_color=CARD_BG)
add_rect(s, 9.1, 3.6, 3.85, 0.05, fill_color=RED)
txb(s, "HEADWINDS", 9.25, 3.67, 3.5, 0.3, font_size=11, bold=True, color=RED)
hw = ["App fatigue among Gen Z", "University budget freezes", "Free competitors (Notion)",
      "Student loan crisis / spend"]
for i, h in enumerate(hw):
    txb(s, f"- {h}", 9.25, 4.05+i*0.33, 3.5, 0.3, font_size=10, color=WHITE)

# Sustainability box
add_rounded_rect(s, 0.4, 5.6, 12.5, 1.3, fill_color=CARD_BG)
add_rect(s, 0.4, 5.6, 12.5, 0.05, fill_color=GREEN)
txb(s, "Sustainability Metrics", 0.6, 5.67, 4, 0.35, font_size=13, bold=True, color=GREEN)

sus_items = [
    ("LTV / CAC Ratio", "36 : 1", ACCENT),
    ("Marginal cost per new user", "~$0", GREEN),
    ("Gross Margin", ">90%", GREEN),
    ("Network effect", "Yes - scales", ACCENT2),
]
for i, (lbl, val, col) in enumerate(sus_items):
    sx = 0.6 + i * 3.0
    txb(s, lbl, sx, 6.1,  2.7, 0.28, font_size=10, color=LIGHT_GRAY)
    txb(s, val, sx, 6.38, 2.7, 0.4,  font_size=18, bold=True, color=col)

txb(s, u"Sources: EdTech CAGR \u2014 Grand View Research 2024  |  Canvas growth \u2014 Instructure 2023 Annual Report  |  LTV/CAC calc: LTV=$4.99x12x3=$179.64; CAC est. $5 (word-of-mouth; cf. Duolingo S-1 $3.20 CAC)  |  Gross margin: hosting $45/mo vs. unlimited users",
    0.4, 6.92, 12.5, 0.28, font_size=8, italic=True, color=LIGHT_GRAY)

# =============================================================================
# SLIDE 11 - COMPETITOR MATRIX
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
slide_number(s, 11)
header_tag(s, "CONNECTIONS  |  SUPPLEMENTAL", color=GREEN)

txb(s, "Competitive Landscape", 0.4, 0.7, 9, 0.7, font_size=42, bold=True, color=WHITE)
txb(s, "StudiBudd is the only tool combining Canvas sync + gamification + student-first design.",
    0.4, 1.35, 11, 0.4, font_size=14, italic=True, color=LIGHT_GRAY)
divider(s, 1.82)

cols_w = [2.8, 1.65, 1.65, 1.65, 1.65, 1.65, 1.65]
cols_x = [0.35]
for cw in cols_w[:-1]:
    cols_x.append(cols_x[-1] + cw)

comp_headers = ["Feature", "StudiBudd", "Notion", "MyStudyLife", "Forest", "Google Cal", "Canvas"]
header_colors= [LIGHT_GRAY, ACCENT, LIGHT_GRAY, LIGHT_GRAY, LIGHT_GRAY, LIGHT_GRAY, LIGHT_GRAY]

for hx, hw, ht, hc in zip(cols_x, cols_w, comp_headers, header_colors):
    bg_c = CARD_BG if ht != "StudiBudd" else RGBColor(0x0D,0x2E,0x44)
    add_rect(s, hx, 2.05, hw, 0.42, fill_color=bg_c)
    txb(s, ht, hx+0.08, 2.1, hw-0.1, 0.35, font_size=12, bold=True,
        color=hc, align=PP_ALIGN.CENTER)

feat_rows = [
    ("Built for students",          True,  False, True,  False, False, True),
    ("Canvas LMS integration",      True,  False, False, False, False, None),
    ("Gamification / motivation",   True,  False, False, True,  False, False),
    ("Assignment auto-import",      True,  False, False, False, False, None),
    ("Course-specific organization",True,  True,  True,  False, False, True),
    ("Customizable aesthetics",     True,  True,  False, False, False, False),
    ("Free tier available",         True,  True,  True,  True,  True,  None),
    ("Price (premium)",             "$4.99", "$10", "Free", "$3.99", "Free", "N/A"),
]

CHECK = "YES"
CROSS = "NO"

for ri, row in enumerate(feat_rows):
    ry = 2.55 + ri * 0.58
    rl_bg = CARD_BG if ri % 2 == 0 else RGBColor(0x13,0x24,0x35)
    for ci, (cx, cw, cell) in enumerate(zip(cols_x, cols_w, row)):
        is_sb = ci == 1
        cell_bg = RGBColor(0x0D,0x2E,0x44) if is_sb else rl_bg
        add_rect(s, cx, ry, cw, 0.5, fill_color=cell_bg)
        if isinstance(cell, bool):
            txt   = CHECK if cell else CROSS
            color = GREEN  if cell else RED
        elif cell is None:
            txt, color = "N/A", LIGHT_GRAY
        else:
            txt   = str(cell)
            color = ACCENT if is_sb else WHITE
        b = is_sb and isinstance(cell, bool) and cell
        txb(s, txt, cx+0.08, ry+0.06, cw-0.12, 0.38,
            font_size=11, bold=b, color=color, align=PP_ALIGN.CENTER)

txb(s, u"StudiBudd's Moat: Only tool at the intersection of Canvas sync + gamification + student-first design.  |  Pricing sources: Notion $10/mo (notion.so), Forest $3.99 (App Store), Google Calendar (free), MyStudyLife (free)",
    0.35, 7.07, 12.5, 0.28, font_size=8, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 12 - CONCLUSION
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s)

# decorative elements
add_rect(s, 0, 0, 0.1, 7.5, fill_color=ACCENT)
add_rect(s, 0.1, 0, 0.05, 7.5, fill_color=ACCENT2)

for r_size, r_col in [(3.5, RGBColor(0x12,0x26,0x38)),
                       (2.5, RGBColor(0x16,0x2E,0x44)),
                       (1.5, RGBColor(0x4A,0xC9,0xFF))]:
    circ = s.shapes.add_shape(9, Inches(12.0-r_size/2), Inches(6.5-r_size/2),
                               Inches(r_size), Inches(r_size))
    circ.fill.solid()
    circ.fill.fore_color.rgb = r_col
    circ.line.fill.background()

txb(s, "The Investment Case", 0.5, 0.5, 10, 0.55,
    font_size=16, color=ACCENT, bold=True)
txb(s, "StudiBudd", 0.5, 1.0, 10, 1.3,
    font_size=72, bold=True, color=WHITE)

add_rect(s, 0.5, 2.2, 6.0, 0.06, fill_color=ACCENT)

# KPI row
kpis = [("NPV", "$91,845"), ("IRR", "254%"), ("Payback", "~10 mo"), ("LTV:CAC", "36:1")]
for i, (k, v) in enumerate(kpis):
    kx = 0.5 + i * 2.8
    txb(s, v, kx, 2.5,  2.6, 0.7, font_size=30, bold=True, color=ACCENT)
    txb(s, k, kx, 3.15, 2.6, 0.4, font_size=12, color=LIGHT_GRAY)

add_rect(s, 0.5, 3.6, 6.0, 0.03, fill_color=ACCENT2)

# bullet summary
summary = [
    "20M underserved students in a $5.2B market — no dominant player",
    "Fully built & deployed — not a concept, not a mockup",
    "Canvas integration live — our structural competitive moat",
    "University-led go-to-market — low CAC, high credibility",
    "Profitable by Year 1, NPV-positive at conservative projections",
]
for i, line in enumerate(summary):
    add_rect(s, 0.5, 3.78+i*0.48, 0.12, 0.28, fill_color=ACCENT)
    txb(s, line, 0.75, 3.75+i*0.48, 7.5, 0.42, font_size=13, color=WHITE)

# right side call to action
add_rounded_rect(s, 8.5, 1.8, 4.5, 3.5, fill_color=CARD_BG)
add_rect(s, 8.5, 1.8, 4.5, 0.06, fill_color=ACCENT2)
txb(s, "Our Ask", 8.7, 1.88, 4.1, 0.4, font_size=16, bold=True, color=ACCENT2)
txb(s, "$10,000 seed investment\nto fund Year 1 marketing\n& university outreach\n\n— OR —\n\nFaculty / administration\nendorsement to launch\nat UNH this semester",
    8.7, 2.38, 4.1, 2.8, font_size=13, color=WHITE)

# unfair advantage box
add_rounded_rect(s, 0.3, 6.6, 12.6, 0.7, fill_color=RGBColor(0x08,0x18,0x28))
add_rect(s, 0.3, 6.6, 0.06, 0.7, fill_color=YELLOW)
txb(s, 'Our Unfair Advantage:  "We are the customer. We built the tool we needed. No outside firm can manufacture that authenticity."',
    0.5, 6.72, 12.2, 0.5, font_size=13, italic=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = r"C:\Users\Owner\Downloads\StudiBudd_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
